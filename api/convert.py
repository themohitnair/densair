import markdown
import re
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_code_box(slide, code_text: str, title: str = None) -> None:
    """Creates a properly formatted code box with monospace font"""
    # Add title if provided
    top = Inches(1.5)
    if title:
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.75)
        )
        title_frame = title_shape.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(0, 0, 139)
        top = Inches(1.5)

    text_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9.0), Inches(5.0))

    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.25)
    text_frame.margin_right = Inches(0.25)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)

    paragraph = text_frame.paragraphs[0]
    paragraph.font.name = "Courier New"
    paragraph.font.size = Pt(12)
    paragraph.text = code_text


def split_code(
    code_text: str, max_chars_per_line: int = 80, max_lines: int = 30
) -> list[str]:
    """Split code into multiple slides if it's too long"""
    lines = code_text.split("\n")
    slides = []
    current_slide = []
    current_chars = 0

    for line in lines:
        if (
            len(current_slide) >= max_lines
            or current_chars + len(line) > max_chars_per_line * max_lines
        ):
            slides.append("\n".join(current_slide))
            current_slide = []
            current_chars = 0

        current_slide.append(line)
        current_chars += len(line)

    if current_slide:
        slides.append("\n".join(current_slide))

    return slides


def create_bullet_slide(
    presentation: Presentation, title: str, bullets: list[tuple[str, int]]
) -> None:
    """Create a slide with properly formatted bullets"""
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)

    if title:
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.word_wrap = True

    for bullet_text, level in bullets:
        p = text_frame.add_paragraph()
        clean_text = re.sub(r"^[-*+•○▪▫]\s*", "", bullet_text.strip())
        p.text = clean_text
        p.level = level
        p.font.size = Pt(18 - (level * 2))
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(12)


def extract_table_data(table_element) -> list[list[str]]:
    """Extract data from HTML table, properly handling headers"""
    table_data = []

    headers = []
    header_row = table_element.find("thead")
    if header_row:
        headers = [th.get_text(strip=True) for th in header_row.find_all(["th", "td"])]
    else:
        first_row = table_element.find("tr")
        if first_row:
            headers = [
                th.get_text(strip=True) for th in first_row.find_all(["th", "td"])
            ]

    if headers:
        table_data.append(headers)

    tbody = table_element.find("tbody")
    if tbody:
        rows = tbody.find_all("tr")
    else:
        rows = (
            table_element.find_all("tr")[1:]
            if headers
            else table_element.find_all("tr")
        )

    for row in rows:
        row_data = [cell.get_text(strip=True) for cell in row.find_all(["td", "th"])]
        if row_data:
            table_data.append(row_data)

    return table_data


def create_table_slide(
    presentation: Presentation, table_data: list[list[str]], title: str = None
) -> None:
    """Create a slide with a properly formatted table"""
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)

    top = Inches(1.0)
    if title:
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
        top = Inches(2.0)

    rows = len(table_data)
    cols = len(table_data[0]) if table_data else 0

    if rows == 0 or cols == 0:
        return

    available_width = Inches(9.0)
    row_height = Inches(0.5)

    # Add table to slide
    table = slide.shapes.add_table(
        rows, cols, Inches(0.5), top, width=available_width, height=row_height * rows
    ).table

    for row_idx, row in enumerate(table_data):
        # Adjust row height
        if row_idx == 0:  # Header row
            table.rows[row_idx].height = Inches(0.6)
        else:
            table.rows[row_idx].height = Inches(0.4)

        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.CENTER

            # Format header row
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(230, 230, 230)
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)


def to_presentation(markdown_input: str) -> Presentation:
    """Convert markdown to a well-formatted PowerPoint presentation"""
    html = markdown.markdown(
        markdown_input, extensions=["fenced_code", "tables", "codehilite", "toc"]
    )

    soup = BeautifulSoup(html, "html.parser")
    presentation = Presentation()

    title = soup.find("h1")
    if title:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = title.text
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

    current_section = []
    current_title = None

    for elem in soup.find_all(["h2", "p", "pre", "ul", "ol", "table"]):
        if elem.name == "h2":
            if current_section:
                if any(item.name == "pre" for item in current_section):
                    # Handle code blocks
                    for code_elem in [
                        item for item in current_section if item.name == "pre"
                    ]:
                        code_text = code_elem.text.strip()
                        code_slides = split_code(code_text)

                        for idx, code_slide in enumerate(code_slides):
                            slide = presentation.slides.add_slide(
                                presentation.slide_layouts[6]
                            )
                            if idx == 0:  # Only add title to first slide of split code
                                create_code_box(slide, code_slide, current_title)
                            else:
                                create_code_box(slide, code_slide)
                else:
                    # Handle bullet points
                    bullets = []
                    for item in current_section:
                        if item.name in ("ul", "ol"):
                            for li in item.find_all("li", recursive=True):
                                level = (
                                    len(list(li.parents)) - len(list(item.parents)) - 1
                                )
                                bullets.append((li.text.strip(), level))

                    if bullets:
                        create_bullet_slide(presentation, current_title, bullets)

            current_section = []
            current_title = elem.text.strip()
        elif elem.name == "table":
            # Handle tables
            table_data = extract_table_data(elem)
            if table_data:
                create_table_slide(presentation, table_data, current_title)
        else:
            current_section.append(elem)

    # Handle last section
    if current_section:
        if any(item.name == "pre" for item in current_section):
            for code_elem in [item for item in current_section if item.name == "pre"]:
                code_text = code_elem.text.strip()
                code_slides = split_code(code_text)

                for idx, code_slide in enumerate(code_slides):
                    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                    if idx == 0:  # Only add title to first slide of split code
                        create_code_box(slide, code_slide, current_title)
                    else:
                        create_code_box(slide, code_slide)
        elif current_section:
            bullets = []
            for item in current_section:
                if item.name in ("ul", "ol"):
                    for li in item.find_all("li", recursive=True):
                        level = len(list(li.parents)) - len(list(item.parents)) - 1
                        bullets.append((li.text.strip(), level))

            if bullets:
                create_bullet_slide(presentation, current_title, bullets)

    return presentation
